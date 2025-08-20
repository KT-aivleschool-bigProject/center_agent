import os
import asyncio
from typing import List, Dict, Any, Optional
from pathlib import Path
import filetype
import magic
from PIL import Image
import pytesseract
import logging
from datetime import datetime
import pandas as pd
import fitz  # PyMuPDF
from pptx import Presentation

from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_openai import ChatOpenAI, OpenAIEmbeddings
from langchain_community.document_loaders import (
    TextLoader,
    PyPDFLoader,
    DirectoryLoader,
    UnstructuredWordDocumentLoader,
    UnstructuredExcelLoader,
    UnstructuredPowerPointLoader
)
from langchain_community.vectorstores import FAISS
from langchain.schema import Document
from langchain.chains import RetrievalQA
from langchain.prompts import PromptTemplate

from langgraph.graph import StateGraph, END
from langgraph.graph.message import add_messages
from typing_extensions import Annotated, TypedDict

# 로깅 설정
logger = logging.getLogger(__name__)


class RAGState(TypedDict):
    """RAG Agent의 상태 정의"""
    question: str
    retrieved_docs: List[Document]
    answer: str
    error: Optional[str]
    metadata: Dict[str, Any]


class RAGAgent:
    """RAG 기반 LangGraph Agent"""
    
    def __init__(self):
        self.name = "RAG Agent"
        self.embeddings = None
        self.vectorstore = None
        self.llm = None
        self.retrieval_qa = None
        self.graph = None
        self.text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=1000,
            chunk_overlap=200,
            length_function=len,
        )
        
        # 문서 저장 경로
        self.docs_path = Path("data/docs")
        self.vector_db_path = Path("data/vector_db")
        
        # 지원하는 파일 형식
        self.supported_extensions = {
            'text': ['.txt', '.md', '.csv'],
            'pdf': ['.pdf'],
            'word': ['.doc', '.docx'],
            'excel': ['.xls', '.xlsx'],
            'powerpoint': ['.ppt', '.pptx'],
            'image': ['.jpg', '.jpeg', '.png', '.gif', '.bmp', '.webp'],
            'hwp': ['.hwp']
        }
        
        # 초기화
        self._initialize_components()
        self._setup_graph()
        
        # 문서는 첫 번째 요청시 로딩 (비동기 함수이므로 별도 호출 필요)
        self._documents_loaded = False
    
    def _initialize_components(self):
        """LangChain 구성 요소 초기화"""
        try:
            # OpenAI API 키 확인
            api_key = os.getenv("OPENAI_API_KEY")
            if not api_key:
                print("  RAG Agent: OpenAI API 키가 설정되지 않았습니다.")
                return
            
            # LLM 초기화
            self.llm = ChatOpenAI(
                model="gpt-4o-mini",
                temperature=0.3,
                api_key=api_key
            )
            
            # 임베딩 모델 초기화
            self.embeddings = OpenAIEmbeddings(api_key=api_key)
            
            # 벡터 저장소 디렉토리 생성
            self.vector_db_path.mkdir(parents=True, exist_ok=True)
            
            # 기존 FAISS 벡터 저장소가 있으면 로드
            try:
                if (self.vector_db_path / "index.faiss").exists():
                    self.vectorstore = FAISS.load_local(
                        str(self.vector_db_path),
                        self.embeddings,
                        allow_dangerous_deserialization=True
                    )
                    self._setup_retrieval_qa()
                    print("  기존 FAISS 벡터 저장소를 로드했습니다.")
                    self._documents_loaded = True
                else:
                    print("  기존 벡터 저장소를 찾을 수 없습니다. 문서를 로딩하겠습니다.")
                    self._documents_loaded = False
            except Exception as e:
                print(f"  FAISS 벡터 저장소 로드 실패: {e}")
                self._documents_loaded = False
            
            print(" RAG Agent 구성 요소가 초기화되었습니다.")
            
        except Exception as e:
            print(f" RAG Agent 초기화 오류: {e}")
    
    def _setup_graph(self):
        """LangGraph 워크플로우 설정"""
        # 상태 그래프 생성
        workflow = StateGraph(RAGState)
        
        # 노드 추가
        workflow.add_node("entry", self._entry_node)
        workflow.add_node("retrieve_docs", self._retrieve_docs_node)
        workflow.add_node("generate_answer", self._generate_answer_node)
        workflow.add_node("output", self._output_node)
        
        # 엣지 설정
        workflow.set_entry_point("entry")
        workflow.add_edge("entry", "retrieve_docs")
        workflow.add_edge("retrieve_docs", "generate_answer")
        workflow.add_edge("generate_answer", "output")
        workflow.add_edge("output", END)
        
        # 그래프 컴파일
        self.graph = workflow.compile()
    
    async def _load_documents(self):
        """문서 디렉토리에서 실제 문서들을 로딩"""
        try:
            # 문서 디렉토리 생성 (존재하지 않을 경우)
            self.docs_path.mkdir(parents=True, exist_ok=True)
            
            # 문서 로딩 및 벡터화
            await self._load_and_vectorize_documents()
            
        except Exception as e:
            print(f" 테스트 문서 로딩 오류: {e}")
    
    async def _load_and_vectorize_documents(self):
        """문서를 로딩하고 벡터화하여 저장"""
        try:
            if not self.embeddings:
                print("  임베딩 모델이 초기화되지 않았습니다.")
                return
            
            documents = []
            
            # 텍스트 파일 직접 로딩
            for txt_file in self.docs_path.glob("*.txt"):
                try:
                    with open(txt_file, 'r', encoding='utf-8') as f:
                        content = f.read()
                        doc = Document(
                            page_content=content,
                            metadata={"source": str(txt_file)}
                        )
                        documents.append(doc)
                        print(f" 로딩 완료: {txt_file.name}")
                except Exception as e:
                    print(f" {txt_file.name} 로딩 실패: {e}")
            
            # PDF 파일 로딩
            for pdf_file in self.docs_path.glob("*.pdf"):
                try:
                    loader = PyPDFLoader(str(pdf_file))
                    docs = loader.load()
                    documents.extend(docs)
                    print(f" PDF 로딩 완료: {pdf_file.name}")
                    print(f" PDF 로딩 완료: {pdf_file.name}")
                except Exception as e:
                    print(f" PDF {pdf_file.name} 로딩 실패: {e}")

            # DOCX 파일 로딩
            for docx_file in self.docs_path.glob("*.docx"):
                try:
                    loader = UnstructuredWordDocumentLoader(str(docx_file))
                    docs = loader.load()
                    documents.extend(docs)
                    print(f" DOCX 로딩 완료: {docx_file.name}")
                except Exception as e:
                    print(f" DOCX {docx_file.name} 로딩 실패: {e}")

            if documents:
                # 문서 분할
                texts = self.text_splitter.split_documents(documents)
                print(f" {len(documents)}개 문서를 {len(texts)}개 청크로 분할했습니다.")
                print(f" {len(documents)}개 문서를 {len(texts)}개 청크로 분할했습니다.")
                
                # 벡터 저장소 생성/업데이트
                if self.vectorstore is None:
                    self.vectorstore = FAISS.from_documents(
                        documents=texts,
                        embedding=self.embeddings
                    )
                    # FAISS 인덱스 저장
                    self.vectorstore.save_local(str(self.vector_db_path))
                    print(" 새로운 FAISS 벡터 저장소를 생성했습니다.")
                else:
                    self.vectorstore.add_documents(texts)
                    # FAISS 인덱스 저장
                    self.vectorstore.save_local(str(self.vector_db_path))
                    print(" 기존 FAISS 벡터 저장소에 문서를 추가했습니다.")
                
                # RetrievalQA 체인 설정
                self._setup_retrieval_qa()
                
                print(f" 벡터화 완료: {len(documents)}개 문서, {len(texts)}개 텍스트 청크")
                print(f" 벡터화 완료: {len(documents)}개 문서, {len(texts)}개 텍스트 청크")
            else:
                print("  로딩할 수 있는 문서가 없습니다.")
                
        except Exception as e:
            print(f" 문서 벡터화 오류: {e}")
            print(f" 문서 벡터화 오류: {e}")
            import traceback
            traceback.print_exc()
    
    def _setup_retrieval_qa(self):
        """RetrievalQA 체인 설정"""
        if not self.vectorstore or not self.llm:
            return
        
        # 프롬프트 템플릿 설정
        prompt_template = """
당신은 문서 기반 질문 답변 AI입니다. 제공된 문서 컨텍스트를 바탕으로 사용자의 질문에 정확하고 도움이 되는 답변을 해주세요.

컨텍스트:
{context}

질문: {question}

답변 가이드라인:
1. 제공된 컨텍스트를 주로 활용하되, 필요시 일반적인 지식도 함께 활용하세요.
2. 답변은 명확하고 구체적으로 작성하세요.
3. 컨텍스트에서 정확한 정보를 찾을 수 없다면 그 사실을 명시하세요.
4. 한국어로 답변하세요.

답변:
"""
        
        PROMPT = PromptTemplate(
            template=prompt_template,
            input_variables=["context", "question"]
        )
        
        # RetrievalQA 체인 생성
        self.retrieval_qa = RetrievalQA.from_chain_type(
            llm=self.llm,
            chain_type="stuff",
            retriever=self.vectorstore.as_retriever(
                search_type="similarity",
                search_kwargs={"k": 4}
            ),
            chain_type_kwargs={"prompt": PROMPT},
            return_source_documents=True
        )
    
    async def _entry_node(self, state: RAGState) -> RAGState:
        """진입점: 질문 입력 처리"""
        print(f"🔍 RAG Agent: 질문 처리 시작 - {state['question']}")
        
        # 코드 관련 질문 분기점 (향후 code_agent 포워딩용)
        question_lower = state["question"].lower()
        code_keywords = ["코드", "구현", "개발", "프로그래밍", "함수", "클래스", "버그"]
        
        if any(keyword in question_lower for keyword in code_keywords):
            state["metadata"] = {
                "potential_code_question": True,
                "forward_to_code_agent": False  # 현재는 비활성화, 향후 확장시 True로 변경
            }
            print("💡 코드 관련 질문으로 감지됨 (향후 code_agent 포워딩 고려)")
        else:
            state["metadata"] = {"potential_code_question": False}
        
        return state
    
    async def _retrieve_docs_node(self, state: RAGState) -> RAGState:
        """문서 검색 노드"""
        try:
            if not self.vectorstore:
                state["error"] = "벡터 저장소가 초기화되지 않았습니다. 문서를 먼저 로딩해주세요."
                print(" 벡터 저장소가 없습니다. 문서 로딩을 다시 시도합니다.")
                await self._load_documents()  # 문서 재로딩 시도
                if not self.vectorstore:
                    return state
            
            # 관련 문서 검색
            retriever = self.vectorstore.as_retriever(
                search_type="similarity",
                search_kwargs={"k": 4}
            )
            
            docs = retriever.get_relevant_documents(state["question"])
            state["retrieved_docs"] = docs
            
            print(f" {len(docs)}개의 관련 문서를 검색했습니다.")
            
            # 검색된 문서가 없을 경우 대체 응답 준비
            if not docs:
                state["error"] = "관련 문서를 찾을 수 없습니다."
                print("  검색된 문서가 없습니다.")
            
        except Exception as e:
            state["error"] = f"문서 검색 중 오류: {str(e)}"
            print(f" 문서 검색 오류: {e}")
        
        return state
    
    async def _generate_answer_node(self, state: RAGState) -> RAGState:
        """답변 생성 노드"""
        try:
            if state.get("error"):
                # 오류가 있지만 문서가 없는 경우 대체 응답 생성
                if "벡터 저장소" in state["error"] or "관련 문서를 찾을 수 없습니다" in state["error"]:
                    state["answer"] = f"""죄송합니다. 현재 문서 저장소에 관련 정보가 없어 정확한 답변을 드릴 수 없습니다.

질문: {state["question"]}

대신 일반적인 도움을 드릴 수 있습니다:
• 프로젝트 문서를 업로드해주시면 더 정확한 답변을 제공할 수 있습니다.
• 구체적인 질문을 해주시면 다른 전문 에이전트가 도움을 드릴 수 있습니다.
• 코드 관련 질문은 'code' 에이전트에게, 문서 작성은 'document' 에이전트에게 문의해보세요."""
                    state["error"] = None  # 에러를 클리어하여 정상 응답으로 처리
                    print(" 대체 응답을 생성했습니다.")
                    return state
                else:
                    return state
            
            if not self.retrieval_qa:
                state["error"] = "RetrievalQA 체인이 초기화되지 않았습니다."
                return state
            
            # RAG 방식으로 답변 생성
            result = self.retrieval_qa({"query": state["question"]})
            state["answer"] = result["result"]
            
            # 소스 문서 정보 추가
            if result.get("source_documents"):
                state["metadata"]["source_count"] = len(result["source_documents"])
                state["answer"] += f"\n\n 참조한 문서: {len(result['source_documents'])}개"
            
            print(" RAG 방식으로 답변이 생성되었습니다.")
            
        except Exception as e:
            state["error"] = f"답변 생성 중 오류: {str(e)}"
            print(f" 답변 생성 오류: {e}")
        
        return state
    
    async def _output_node(self, state: RAGState) -> RAGState:
        """출력 노드"""
        if state.get("error"):
            state["answer"] = f"처리 중 오류가 발생했습니다: {state['error']}"
        
        print(" RAG Agent 처리 완료")
        return state
    
    async def process(self, message: str) -> str:
        """메시지 처리 (FastAPI 통합용)"""
        try:
            # 첫 번째 요청시 문서 로딩
            if not self._documents_loaded:
                await self._load_documents()
                self._documents_loaded = True
            
            # 초기 상태 설정
            initial_state = RAGState(
                question=message,
                retrieved_docs=[],
                answer="",
                error=None,
                metadata={}
            )
            
            # 그래프 실행
            if self.graph:
                result = await self.graph.ainvoke(initial_state)
                return result["answer"]
            else:
                return "RAG Agent가 초기화되지 않았습니다. OpenAI API 키를 확인해주세요."
        
        except Exception as e:
            return f"RAG Agent 처리 중 오류가 발생했습니다: {str(e)}"
    
    async def add_documents(self, file_paths: List[str]) -> str:
        """새 문서 추가"""
        try:
            documents = []
            for file_path in file_paths:
                path = Path(file_path)
                if path.suffix == ".txt":
                    loader = TextLoader(str(path))
                elif path.suffix == ".pdf":
                    loader = PyPDFLoader(str(path))
                elif path.suffix == ".docx":
                    loader = UnstructuredWordDocumentLoader(str(path))
                else:
                    continue
                
                docs = loader.load()
                documents.extend(docs)
            
            if documents and self.vectorstore:
                texts = self.text_splitter.split_documents(documents)
                self.vectorstore.add_documents(texts)
                return f"{len(documents)}개 문서가 추가되었습니다."
            else:
                return "추가할 수 있는 문서가 없습니다."
        
        except Exception as e:
            return f"문서 추가 중 오류: {str(e)}"

    def _determine_file_type(self, file_path: str) -> str:
        """파일 형식 판별"""
        path = Path(file_path)
        extension = path.suffix.lower()
        
        for file_type, extensions in self.supported_extensions.items():
            if extension in extensions:
                return file_type
        
        # MIME 타입으로 재확인
        try:
            mime_type = magic.from_file(file_path, mime=True)
            if mime_type.startswith('text/'):
                return 'text'
            elif mime_type == 'application/pdf':
                return 'pdf'
            elif 'word' in mime_type or 'document' in mime_type:
                return 'word'
            elif 'excel' in mime_type or 'spreadsheet' in mime_type:
                return 'excel'
            elif 'powerpoint' in mime_type or 'presentation' in mime_type:
                return 'powerpoint'
            elif mime_type.startswith('image/'):
                return 'image'
        except:
            pass
        
        return 'unknown'

    def _extract_text_from_image(self, file_path: str) -> str:
        """이미지에서 OCR로 텍스트 추출"""
        try:
            image = Image.open(file_path)
            text = pytesseract.image_to_string(image, lang='kor+eng')
            return text.strip()
        except Exception as e:
            print(f"이미지 OCR 추출 실패 ({file_path}): {e}")
            return ""

    def _load_single_document(self, file_path: str, project_id: str = None, file_id: str = None) -> List[Document]:
        """단일 파일을 로드하여 Document 객체 리스트 반환"""
        documents = []
        path = Path(file_path)
        file_type = self._determine_file_type(file_path)
        
        metadata = {
            "source": str(path),
            "filename": path.name,
            "file_type": file_type,
        }
        
        if project_id:
            metadata["project_id"] = project_id
        if file_id:
            metadata["file_id"] = file_id
        
        try:
            if file_type == 'text':
                with open(file_path, 'r', encoding='utf-8') as f:
                    content = f.read()
                documents.append(Document(page_content=content, metadata=metadata))
                
            elif file_type == 'pdf':
                loader = PyPDFLoader(file_path)
                docs = loader.load()
                for doc in docs:
                    doc.metadata.update(metadata)
                documents.extend(docs)
                
            elif file_type == 'word':
                loader = UnstructuredWordDocumentLoader(file_path)
                docs = loader.load()
                for doc in docs:
                    doc.metadata.update(metadata)
                documents.extend(docs)
                
            elif file_type == 'excel':
                loader = UnstructuredExcelLoader(file_path)
                docs = loader.load()
                for doc in docs:
                    doc.metadata.update(metadata)
                documents.extend(docs)
                
            elif file_type == 'powerpoint':
                loader = UnstructuredPowerPointLoader(file_path)
                docs = loader.load()
                for doc in docs:
                    doc.metadata.update(metadata)
                documents.extend(docs)
                
            elif file_type == 'image':
                text = self._extract_text_from_image(file_path)
                if text:
                    documents.append(Document(
                        page_content=f"[이미지에서 추출된 텍스트]\n{text}",
                        metadata=metadata
                    ))
                    
            elif file_type == 'hwp':
                # HWP 파일은 별도 처리 필요 (현재는 스킵)
                print(f"HWP 파일은 현재 지원되지 않습니다: {file_path}")
                
            else:
                print(f"지원되지 않는 파일 형식: {file_path}")
                
        except Exception as e:
            print(f"파일 로딩 실패 ({file_path}): {e}")
            
        return documents

    def load_project_documents(self, project_id: str = None) -> List[Document]:
        """ai-server/data/docs 경로의 모든 파일을 로드"""
        documents = []
        
        # 디렉토리가 존재하지 않으면 생성
        self.docs_path.mkdir(parents=True, exist_ok=True)
        
        # 지원하는 모든 확장자 패턴
        all_extensions = []
        for extensions in self.supported_extensions.values():
            all_extensions.extend(extensions)
        
        for extension in all_extensions:
            for file_path in self.docs_path.glob(f"*{extension}"):
                docs = self._load_single_document(str(file_path), project_id)
                documents.extend(docs)
                
        return documents

    async def process_new_document(self, file_path: str, project_id: str = None, file_id: str = None) -> str:
        """새로 업로드된 단일 문서를 처리"""
        try:
            logger.info(f"문서 처리 시작: {file_path}")
            
            # 파일 존재 확인
            if not Path(file_path).exists():
                return f"파일을 찾을 수 없습니다: {file_path}"
            
            # 벡터 스토어가 초기화되지 않았다면 초기화
            if not self.vectorstore:
                await self._load_documents()
            
            # 파일 내용 읽기
            content = self.load_document_content(file_path)
            
            if not content.strip():
                return f"파일에서 내용을 추출할 수 없습니다: {file_path}"
            
            # 문서 메타데이터 생성
            metadata = {
                "filename": os.path.basename(file_path),
                "project_id": project_id,
                "file_id": file_id,
                "file_path": file_path,
                "upload_time": datetime.now().isoformat(),
                "file_type": self.get_file_extension(file_path),
                "content_length": len(content)
            }
            
            # 벡터 스토어에 추가
            await self.add_document_to_vector_store(content, metadata)
            
            # RetrievalQA 재설정
            self._setup_retrieval_qa()
            
            result_msg = f"문서 처리 완료: {metadata['filename']} (프로젝트: {project_id}, 내용 길이: {len(content)}자)"
            logger.info(result_msg)
            return result_msg
            
        except Exception as e:
            error_msg = f"문서 처리 중 오류: {str(e)}"
            logger.error(error_msg)
            import traceback
            traceback.print_exc()
            return error_msg

    def _check_document_exists(self, file_name: str, project_id: str = None) -> bool:
        """문서가 이미 존재하는지 확인"""
        # 실제 구현에서는 벡터 스토어의 메타데이터를 검색해야 함
        # 현재는 단순히 False 반환 (항상 새 문서로 처리)
        return False

    async def _remove_existing_documents(self, file_name: str, project_id: str = None):
        """기존 문서를 벡터 스토어에서 제거"""
        # FAISS는 직접적인 문서 제거가 어려우므로
        # 실제 구현에서는 ChromaDB 등 다른 벡터 DB 사용 권장
        pass

    async def search_documents(self, query: str, project_id: str = None, limit: int = 5) -> List[Dict]:
        """문서 검색 (프로젝트 필터링 지원)"""
        try:
            if not self.vectorstore:
                return []
            
            # 벡터 검색 수행
            all_docs = self.vectorstore.similarity_search(query, k=limit*2)
            
            # 프로젝트 ID로 필터링 (지정된 경우)
            if project_id:
                filtered_docs = [
                    doc for doc in all_docs 
                    if doc.metadata.get('project_id') == project_id
                ]
                docs = filtered_docs[:limit]
            else:
                docs = all_docs[:limit]
            
            # 결과 포맷팅
            results = []
            for doc in docs:
                results.append({
                    "content": doc.page_content[:500] + "..." if len(doc.page_content) > 500 else doc.page_content,
                    "metadata": doc.metadata,
                    "score": getattr(doc, 'score', 0.0)
                })
            
            logger.info(f"문서 검색 완료: 쿼리='{query}', 프로젝트={project_id}, 결과={len(results)}개")
            return results
            
        except Exception as e:
            logger.error(f"문서 검색 실패: {e}")
            return []

    def get_file_extension(self, file_path: str) -> str:
        """파일 확장자 추출"""
        return Path(file_path).suffix.lower().lstrip('.')

    def load_document_content(self, file_path: str) -> str:
        """파일 형식에 따라 내용 추출"""
        file_ext = self.get_file_extension(file_path)
        
        try:
            if file_ext == 'txt':
                return self.load_text_file(file_path)
            elif file_ext == 'pdf':
                return self.load_pdf_file(file_path)
            elif file_ext in ['doc', 'docx']:
                return self.load_word_file(file_path)
            elif file_ext in ['xls', 'xlsx']:
                return self.load_excel_file(file_path)
            elif file_ext in ['ppt', 'pptx']:
                return self.load_powerpoint_file(file_path)
            elif file_ext in ['jpg', 'jpeg', 'png', 'gif', 'bmp', 'webp']:
                return self.extract_text_from_image(file_path)
            elif file_ext == 'hwp':
                logger.warning(f"HWP 파일은 현재 지원되지 않습니다: {file_path}")
                return ""
            else:
                logger.warning(f"지원하지 않는 파일 형식: {file_ext}")
                return ""
        except Exception as e:
            logger.error(f"파일 내용 추출 실패: {file_path}, 오류: {str(e)}")
            return ""

    def load_text_file(self, file_path: str) -> str:
        """텍스트 파일 로드"""
        with open(file_path, 'r', encoding='utf-8') as f:
            return f.read()

    def load_pdf_file(self, file_path: str) -> str:
        """PDF 파일 로드 (PyMuPDF 사용)"""
        try:
            doc = fitz.open(file_path)
            text = ""
            for page in doc:
                text += page.get_text()
            doc.close()
            return text
        except Exception as e:
            logger.error(f"PDF 파일 로드 실패: {e}")
            # 대안으로 PyPDFLoader 사용
            try:
                loader = PyPDFLoader(file_path)
                pages = loader.load()
                return "\n".join([page.page_content for page in pages])
            except Exception as e2:
                logger.error(f"PyPDFLoader도 실패: {e2}")
                return ""

    def load_word_file(self, file_path: str) -> str:
        """Word 문서 로드"""
        try:
            loader = UnstructuredWordDocumentLoader(file_path)
            docs = loader.load()
            return "\n".join([doc.page_content for doc in docs])
        except Exception as e:
            logger.error(f"Word 파일 로드 실패: {e}")
            return ""

    def load_excel_file(self, file_path: str) -> str:
        """Excel 파일 로드"""
        try:
            # pandas로 읽기
            df = pd.read_excel(file_path, sheet_name=None)  # 모든 시트 읽기
            content = ""
            for sheet_name, sheet_df in df.items():
                content += f"\n=== 시트: {sheet_name} ===\n"
                content += sheet_df.to_string(index=False)
                content += "\n"
            return content
        except Exception as e:
            logger.error(f"Excel 파일 로드 실패: {e}")
            try:
                # 대안으로 UnstructuredExcelLoader 사용
                loader = UnstructuredExcelLoader(file_path)
                docs = loader.load()
                return "\n".join([doc.page_content for doc in docs])
            except Exception as e2:
                logger.error(f"UnstructuredExcelLoader도 실패: {e2}")
                return ""

    def load_powerpoint_file(self, file_path: str) -> str:
        """PowerPoint 파일 로드"""
        try:
            prs = Presentation(file_path)
            content = ""
            for i, slide in enumerate(prs.slides):
                content += f"\n=== 슬라이드 {i+1} ===\n"
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        content += shape.text + "\n"
            return content
        except Exception as e:
            logger.error(f"PowerPoint 파일 로드 실패: {e}")
            try:
                # 대안으로 UnstructuredPowerPointLoader 사용
                loader = UnstructuredPowerPointLoader(file_path)
                docs = loader.load()
                return "\n".join([doc.page_content for doc in docs])
            except Exception as e2:
                logger.error(f"UnstructuredPowerPointLoader도 실패: {e2}")
                return ""

    def extract_text_from_image(self, file_path: str) -> str:
        """이미지에서 OCR로 텍스트 추출"""
        try:
            image = Image.open(file_path)
            text = pytesseract.image_to_string(image, lang='kor+eng')
            return f"[이미지에서 추출된 텍스트]\n{text.strip()}"
        except Exception as e:
            logger.error(f"이미지 OCR 추출 실패 ({file_path}): {e}")
            return ""

    async def add_document_to_vector_store(self, content: str, metadata: dict):
        """문서를 벡터 스토어에 추가"""
        try:
            if not content.strip():
                logger.warning(f"빈 문서 내용: {metadata.get('filename', 'unknown')}")
                return
            
            # 기존 문서 중복 체크
            filename = metadata.get('filename', '')
            project_id = metadata.get('project_id', '')
            
            existing_doc = self.check_document_exists(filename, project_id)
            if existing_doc:
                logger.info(f"기존 문서 업데이트: {filename}")
                await self.update_existing_document(content, metadata)
            else:
                logger.info(f"새 문서 추가: {filename}")
                await self.add_new_document(content, metadata)
                
        except Exception as e:
            logger.error(f"벡터 스토어 업데이트 실패: {str(e)}")
            raise

    def check_document_exists(self, filename: str, project_id: str = None) -> bool:
        """문서가 이미 존재하는지 확인"""
        # 현재 FAISS는 메타데이터 기반 검색이 제한적이므로 단순히 False 반환
        # 실제 구현에서는 ChromaDB 등 메타데이터 필터링이 가능한 벡터DB 사용 권장
        return False

    async def update_existing_document(self, content: str, metadata: dict):
        """기존 문서 업데이트"""
        # FAISS는 직접적인 업데이트가 어려우므로 새 문서로 추가
        await self.add_new_document(content, metadata)

    async def add_new_document(self, content: str, metadata: dict):
        """새 문서 추가"""
        try:
            # Document 객체 생성
            doc = Document(page_content=content, metadata=metadata)
            
            # 문서 분할
            texts = self.text_splitter.split_documents([doc])
            
            # 벡터 스토어에 추가
            if self.vectorstore:
                self.vectorstore.add_documents(texts)
            else:
                # 벡터 스토어가 없으면 새로 생성
                self.vectorstore = FAISS.from_documents(
                    documents=texts,
                    embedding=self.embeddings
                )
            
            # 저장
            self.vectorstore.save_local(str(self.vector_db_path))
            logger.info(f"문서 벡터화 완료: {len(texts)}개 청크")
            
        except Exception as e:
            logger.error(f"새 문서 추가 실패: {e}")
            raise
